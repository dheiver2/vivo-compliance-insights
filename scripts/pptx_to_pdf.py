"""
Renderiza cada slide do PPTX como imagem PNG e combina em um PDF.
Usa: python-pptx, Pillow
"""
import io
import re
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

PPTX_PATH = Path("/Users/dheiver/Documents/vivo-compliance-insights-1/VoiceAudit-Apresentacao-Vendas.pptx")
PDF_PATH  = Path("/Users/dheiver/Documents/vivo-compliance-insights-1/VoiceAudit-Apresentacao-Vendas.pdf")

W_PX, H_PX = 1280, 720   # output resolution
SCALE = W_PX / 9144000   # EMU → pixel  (9144000 EMU = 10 inches @ 914400 EMU/in)

def emu_to_px(v):
    return int(v * SCALE)

def rgb_from_pptx(color_obj):
    """Return (r,g,b) tuple or None."""
    try:
        if color_obj and color_obj.type is not None:
            rgb = color_obj.rgb
            return (rgb.red, rgb.green, rgb.blue)
    except Exception:
        pass
    return None

def hex_to_rgb(h):
    h = h.lstrip("#")
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

def slide_bg_color(slide, prs):
    """Try to get solid background fill color."""
    try:
        bg = slide.background.fill
        if bg.fore_color and bg.fore_color.rgb:
            c = bg.fore_color.rgb
            return (c.red, c.green, c.blue)
    except Exception:
        pass
    # fallback: slide layout / master bg
    try:
        bg = slide.slide_layout.background.fill
        if bg.fore_color and bg.fore_color.rgb:
            c = bg.fore_color.rgb
            return (c.red, c.green, c.blue)
    except Exception:
        pass
    return (255, 255, 255)

def render_slide(slide, prs):
    img = Image.new("RGB", (W_PX, H_PX), slide_bg_color(slide, prs))
    draw = ImageDraw.Draw(img)

    for shape in slide.shapes:
        x = emu_to_px(shape.left  or 0)
        y = emu_to_px(shape.top   or 0)
        w = emu_to_px(shape.width or 0)
        h = emu_to_px(shape.height or 0)

        # ── PICTURES ──────────────────────────────────────
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                blob = shape.image.blob
                pic  = Image.open(io.BytesIO(blob)).convert("RGB")
                pic  = pic.resize((max(w,1), max(h,1)), Image.LANCZOS)
                img.paste(pic, (x, y))
            except Exception:
                pass
            continue

        # ── SHAPES (rect / oval / line) ────────────────────
        if shape.shape_type in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM,
                                 1,  # MSO_SHAPE_TYPE.RECTANGLE
                                 ):
            fill_color = None
            try:
                fill = shape.fill
                if fill.type is not None:
                    fill_color = rgb_from_pptx(fill.fore_color)
            except Exception:
                pass

            if fill_color:
                # Use shape name hint for oval vs rect
                sn = shape.shape_type
                name = (shape.name or "").lower()
                if "oval" in name or "ellipse" in name or "circle" in name:
                    draw.ellipse([x, y, x+w, y+h], fill=fill_color)
                else:
                    draw.rectangle([x, y, x+w, y+h], fill=fill_color)

        # ── TEXT FRAMES ────────────────────────────────────
        if shape.has_text_frame:
            tf = shape.text_frame
            # Collect all runs with their formatting
            para_y = y + 4
            for para in tf.paragraphs:
                line_parts = []
                for run in para.runs:
                    txt = run.text
                    if not txt.strip():
                        continue
                    # font size
                    fs = 14
                    try:
                        fs = int((run.font.size or para.runs[0].font.size or Pt(14)) / 12700)
                    except Exception:
                        pass
                    fs = max(8, min(fs, 80))
                    # color
                    color = (255, 255, 255)
                    try:
                        c = run.font.color.rgb
                        color = (c.red, c.green, c.blue)
                    except Exception:
                        pass
                    bold = False
                    try:
                        bold = bool(run.font.bold)
                    except Exception:
                        pass
                    line_parts.append((txt, fs, color, bold))

                if not line_parts and para.text.strip():
                    # fallback: whole paragraph
                    txt = para.text
                    fs = 14
                    try:
                        fs = int((para.runs[0].font.size or Pt(14)) / 12700) if para.runs else 14
                    except Exception:
                        pass
                    fs = max(8, min(fs, 80))
                    line_parts = [(txt, fs, (200,200,200), False)]

                if not line_parts:
                    para_y += 4
                    continue

                # Use the max font size in the line for height
                max_fs = max(p[1] for p in line_parts)
                try:
                    font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", max_fs)
                except Exception:
                    font = ImageFont.load_default()

                cur_x = x + 4
                for txt, fs, color, bold in line_parts:
                    try:
                        f = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", fs)
                    except Exception:
                        f = ImageFont.load_default()
                    draw.text((cur_x, para_y), txt, font=f, fill=color)
                    bbox = draw.textbbox((0,0), txt, font=f)
                    cur_x += (bbox[2] - bbox[0]) + 2

                para_y += max_fs + 3

    return img

def main():
    prs = Presentation(str(PPTX_PATH))
    images = []
    for i, slide in enumerate(prs.slides):
        print(f"  Renderizando slide {i+1}/{len(prs.slides)}…")
        img = render_slide(slide, prs)
        images.append(img)

    print(f"\nCombinando {len(images)} slides em PDF…")
    images[0].save(
        str(PDF_PATH),
        save_all=True,
        append_images=images[1:],
        resolution=128
    )
    print(f"✅ PDF gerado: {PDF_PATH}  ({len(images)} páginas)")

if __name__ == "__main__":
    main()
